#!/usr/bin/env python3
"""Build TDA seminar PowerPoint from slide script."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor

# Slide content: (title, bullets list)
SLIDES = [
    ("Feedback Loops as Loops: Topological Data Analysis of Genetic Regulatory Circuits",
     ["Gary Welz | CopernicusAI / CUNY Graduate Center (PoI)",
      "Collaboration with Dr. Mikael Vejdemo-Johansson",
      "March 13, 2026"]),
    ("The Innovation: Text to Visual Data",
     ["Traditional TDA: Numerical data -> Feature vectors -> Topology",
      "This work: Text (papers) -> Visual flowcharts -> Features -> Topology",
      "Key: Mermaid Markdown converts textual process descriptions into structured flowcharts",
      "Flowcharts become visual data; TDA reveals structure",
      "We extract topology from descriptions, not from direct measurements"]),
    ("The Question",
     ["Can we detect regulatory structure (feedback, cascades) from circuit topology?",
      "Feedback loops are literally loops -- they should appear in H1",
      "Can text-derived visual data support that?"]),
    ("The GLMP Database",
     ["108 genetic regulatory circuits",
      "E. coli (66), S. cerevisiae (38), Bacillus subtilis (4)",
      "Each process: Mermaid flowchart (nodes, conditionals, OR/AND/NOT gates)",
      "Examples: lac operon, SOS response, two-component signaling",
      "Data: github.com/garywelz/glmp"]),
    ("From Flowcharts to Features",
     ["Node count, conditional count",
      "OR gates, AND gates, NOT gates",
      "Standardized (zero mean, unit variance)",
      "108 processes x 5 features",
      "These capture circuit complexity and logic structure"]),
    ("TDA Pipeline",
     ["Feature matrix -> Euclidean distance",
      "Vietoris-Rips filtration",
      "Ripser (maxdim=2), cocycle extraction",
      "Output: persistence diagrams (H0, H1, H2)",
      "Cocycles tell us which processes form each H1 loop"]),
    ("Persistence Diagram",
     ["H0: 108 components (one per process)",
      "H1: 33 loops",
      "H2: 4 voids",
      "Key question: Do these loops correspond to biological structure?",
      "[Insert image: glmp_persistence_diagram.png]"]),
    ("Top H1 Loop #1 (Persistence = 0.330)",
     ["27 processes: ara operon, SOS, stringent response, catabolite repression, Pho regulon, quorum sensing, heat shock, GAL, MAPK mating",
      "Common thread: regulatory circuits with feedback",
      "Cross-organism: E. coli, yeast, Bacillus"]),
    ("Top H1 Loop #3 (Persistence = 0.302)",
     ["4 processes: Lac operon, antibiotic efflux pumps, phosphate regulation, translation termination",
      "Lac operon = textbook negative feedback",
      "Topology groups by regulatory logic, not just metabolic function"]),
    ("Top H1 Loop #5 (Persistence = 0.231)",
     ["3 processes: Two-component EnvZ-OmpR, oxidative stress response, yeast ER-associated degradation",
      "EnvZ-OmpR = paradigm two-component signaling with feedback",
      "Cross-organism: regulatory logic transcends species"]),
    ("Biological Coherence Check",
     ["Known feedback circuits: lac, trp, ara operons; two-component; SOS, heat shock; catabolite repression, Pho",
      "Result: lac (Loop #3), two-component (Loop #5), many stress/nutrient circuits (Loop #1)",
      "Interpretation: topology captures genuine regulatory architecture"]),
    ("Organism Patterns",
     ["Species-specific: Loop #2 (all yeast), Loop #3 (all E. coli)",
      "Cross-organism: Loop #1, Loop #5",
      "Some motifs are universal; others organism-specific"]),
    ("Limitations and Caveats",
     ["Sample size: 108 (small but meaningful)",
      "Feature choice: structural counts only; could add graph metrics",
      "Biological validation: domain experts should validate",
      "LLM-generated flowcharts: fact-checking needed",
      "Open question: Does topology predict function or correlate with known biology?"]),
    ("Next Steps",
     ["Mapper: circuit classes as nodes",
      "Persistent cohomology: circular coordinates for feedback depth",
      "Scaling: 200-500+ processes; 314 across biology, chemistry, physics, math, CS",
      "Collaboration: Dr. Vejdemo-Johansson, CUNY TDA group; biologist validation",
      "Open source: github.com/garywelz/glmp/tree/main/tda-analysis"]),
    ("Acknowledgments and Questions",
     ["Collaborators: Dr. Mikael Vejdemo-Johansson, Jordan Matuszewski",
      "GLMP: github.com/garywelz/glmp",
      "TDA analysis: github.com/garywelz/glmp/tree/main/tda-analysis",
      "Questions?",
      "",
      "Gary Welz | garywelz@gmail.com | 917-593-2537"]),
]


def add_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(6)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for title, bullets in SLIDES:
        add_slide(prs, title, bullets)
    out_path = "TDA_Seminar_Slides.pptx"
    prs.save(out_path)
    print(f"Saved {out_path}")


if __name__ == "__main__":
    main()
